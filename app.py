import streamlit as st
from pathlib import Path
from pptx import Presentation
import fitz  # PyMuPDF
import pandas as pd
from openai import OpenAI
import numpy as np
from scipy.spatial.distance import cosine
import os
import json

# Initialize OpenAI client
client = OpenAI(api_key="sk-your-actual-key-here")  # Replace with your actual OpenAI key

# File to store embeddings cache
CACHE_FILE = "embeddings_cache.json"

@st.cache_data
def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_content.append(" ".join(slide_text))
        
        return " ".join(text_content)
    except Exception as e:
        st.warning(f"Could not process {file_path}: {str(e)}")
        return ""

@st.cache_data
def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    try:
        doc = fitz.open(file_path)
        text_content = []
        
        for page in doc:
            text_content.append(page.get_text())
        
        return " ".join(text_content)
    except Exception as e:
        st.warning(f"Could not process {file_path}: {str(e)}")
        return ""

def get_embedding(text, client):
    """Get embeddings for a text using OpenAI."""
    try:
        response = client.embeddings.create(
            input=text[:8000],  # Limit text length to prevent token overflow
            model="text-embedding-ada-002"
        )
        return response.data[0].embedding
    except Exception as e:
        st.error(f"Error getting embedding: {str(e)}")
        return None

def load_or_create_embeddings_cache():
    """Load existing embeddings cache or create new one."""
    if os.path.exists(CACHE_FILE):
        with open(CACHE_FILE, 'r') as f:
            return json.load(f)
    return {}

def save_embeddings_cache(cache):
    """Save embeddings cache to file."""
    with open(CACHE_FILE, 'w') as f:
        json.dump(cache, f)

def process_presentations_folder(folder_path):
    """Process all presentations in the folder and return their embeddings."""
    cache = load_or_create_embeddings_cache()
    presentations_data = []
    
    for file_path in Path(folder_path).rglob('*'):
        if file_path.suffix.lower() in ['.pptx', '.ppt', '.pdf']:
            file_str = str(file_path)
            
            # Check if we have cached embeddings
            if file_str in cache:
                presentations_data.append({
                    'path': file_str,
                    'embedding': cache[file_str]['embedding'],
                    'content': cache[file_str]['content']
                })
                continue
            
            # Extract text based on file type
            if file_path.suffix.lower() in ['.pptx', '.ppt']:
                content = extract_text_from_pptx(file_path)
            else:
                content = extract_text_from_pdf(file_path)
            
            # Get embedding for the content
            if content.strip():
                embedding = get_embedding(content, client)
                if embedding:
                    # Cache the results
                    cache[file_str] = {
                        'embedding': embedding,
                        'content': content
                    }
                    presentations_data.append({
                        'path': file_str,
                        'embedding': embedding,
                        'content': content
                    })
    
    # Save updated cache
    save_embeddings_cache(cache)
    return presentations_data

def search_presentations(query, presentations_data, top_k=5):
    """Search through presentations using semantic similarity."""
    if not presentations_data:
        return []
    
    # Get query embedding
    query_embedding = get_embedding(query, client)
    if not query_embedding:
        return []
    
    # Calculate similarities
    similarities = []
    for pres in presentations_data:
        similarity = 1 - cosine(query_embedding, pres['embedding'])
        similarities.append({
            'path': pres['path'],
            'similarity': similarity,
            'content': pres['content']
        })
    
    # Sort by similarity and return top results
    similarities.sort(key=lambda x: x['similarity'], reverse=True)
    return similarities[:top_k]

def main():
    st.set_page_config(
        page_title="Presentation Content Search",
        page_icon="🔍",
        layout="wide"
    )
    
    st.title("🔍 Presentation Content Search")
    st.write("Search through the presentation repository")

    # Sidebar for configuration
    with st.sidebar:
        st.header("Settings")
        presentations_dir = st.text_input(
            "Presentations Directory",
            value="presentations",
            help="Path to the directory containing presentations"
        )
        
        num_results = st.slider(
            "Number of results to show",
            min_value=1,
            max_value=20,
            value=5
        )
        
        min_similarity = st.slider(
            "Minimum similarity score",
            min_value=0.0,
            max_value=1.0,
            value=0.5,
            step=0.05
        )

    # Main search interface
    query = st.text_area(
        "What are you looking for?",
        placeholder="Example: presentations about climate change impacts and mitigation strategies",
        height=100
    )

    search_button = st.button("🔍 Search", type="primary")

    if search_button and query:
        with st.spinner("Searching through presentations..."):
            # Process presentations folder
            presentations_data = process_presentations_folder(presentations_dir)
            
            if not presentations_data:
                st.error(f"No presentations found in {presentations_dir}")
                return
            
            # Search presentations
            results = search_presentations(
                query,
                presentations_data,
                top_k=num_results
            )
            
            if results:
                st.subheader(f"Found {len(results)} relevant presentations")
                
                for i, result in enumerate(results, 1):
                    similarity_percentage = result['similarity'] * 100
                    if similarity_percentage >= min_similarity * 100:
                        with st.expander(f"📊 {Path(result['path']).name}", expanded=i==1):
                            cols = st.columns([3, 1])
                            with cols[0]:
                                st.text("File path: " + result['path'])
                                st.write("**Content Preview:**")
                                # Show first 500 characters of content with "..." if truncated
                                preview = result['content'][:500]
                                if len(result['content']) > 500:
                                    preview += "..."
                                st.write(preview)
                            with cols[1]:
                                st.metric("Relevance", f"{similarity_percentage:.1f}%")
                                file_type = Path(result['path']).suffix.upper()[1:]
                                st.write(f"**File type:** {file_type}")
                
                # Download results
                results_df = pd.DataFrame([{
                    'filename': Path(r['path']).name,
                    'path': r['path'],
                    'similarity': f"{r['similarity']*100:.1f}%",
                } for r in results])
                
                st.download_button(
                    "Download Results as CSV",
                    results_df.to_csv(index=False),
                    "search_results.csv",
                    "text/csv",
                    key='download-csv'
                )
            else:
                st.warning("No matching presentations found. Try adjusting your search terms.")

if __name__ == "__main__":
    main()
